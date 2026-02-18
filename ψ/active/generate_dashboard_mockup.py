"""
Generate Core Dashboard Mockup — reference สำหรับทีม Vibe Code
"""

import sys
sys.stdout.reconfigure(encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
OUTPUT_PATH = os.path.join(SCRIPT_DIR, "CoreDashboard-Mockup.pptx")

# Colors
BG = RGBColor(0x1a, 0x1a, 0x2e)
SIDEBAR_BG = RGBColor(0x12, 0x12, 0x22)
CARD_BG = RGBColor(0x22, 0x22, 0x3a)
GOLD = RGBColor(0xe5, 0xc1, 0x00)
ORANGE = RGBColor(0xf3, 0x9c, 0x12)
BLUE = RGBColor(0x34, 0x98, 0xdb)
GREEN = RGBColor(0x27, 0xae, 0x60)
RED = RGBColor(0xe7, 0x4c, 0x3c)
PURPLE = RGBColor(0x9b, 0x59, 0xb6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xBD, 0xC3, 0xC7)
DIM = RGBColor(0x7f, 0x8c, 0x8d)
DARK_CARD = RGBColor(0x16, 0x16, 0x2a)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def bg(slide, color=BG):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def box(slide, l, t, w, h, fill=CARD_BG, border=None, bw=1.5, radius=True):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    s = slide.shapes.add_shape(shape_type, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if border:
        s.line.color.rgb = border
        s.line.width = Pt(bw)
    else:
        s.line.fill.background()
    return s


def txt(slide, l, t, w, h, text, sz=18, c=WHITE, bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.color.rgb = c
    p.font.bold = bold
    p.font.name = 'Tahoma'
    p.alignment = align
    return tb


def bullets(slide, l, t, w, h, items, sz=14, c=LIGHT, sp=6):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(sz)
        p.font.color.rgb = c
        p.font.name = 'Tahoma'
        p.space_after = Pt(sp)
    return tb


# ============================================================
# SLIDE 1 — Overview: Dashboard Layout
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
txt(s, Inches(0.5), Inches(0.2), Inches(12), Inches(0.6),
    "Core Dashboard — Mockup Reference", sz=28, c=WHITE, bold=True)
txt(s, Inches(0.5), Inches(0.7), Inches(12), Inches(0.4),
    "สำหรับทีม Vibe Code  |  Dark theme  |  Control Center style", sz=14, c=GOLD)

# === SIDEBAR ===
box(s, Inches(0.3), Inches(1.2), Inches(1.8), Inches(5.9), fill=SIDEBAR_BG, border=RGBColor(0x2a, 0x2a, 0x40))

txt(s, Inches(0.4), Inches(1.3), Inches(1.6), Inches(0.4),
    "SWT", sz=16, c=GOLD, bold=True, align=PP_ALIGN.CENTER)
txt(s, Inches(0.4), Inches(1.7), Inches(1.6), Inches(0.3),
    "Core Dashboard", sz=9, c=DIM, align=PP_ALIGN.CENTER)

menu_items = [
    ("Overview", GOLD),
    ("Users", LIGHT),
    ("Roles", LIGHT),
    ("Report", LIGHT),
    ("Alert", LIGHT),
    ("Kiosk", LIGHT),
    ("Settings", DIM),
]
for i, (label, color) in enumerate(menu_items):
    y = Inches(2.2 + i * 0.45)
    if i == 0:
        box(s, Inches(0.4), y - Inches(0.05), Inches(1.5), Inches(0.38),
            fill=RGBColor(0x2a, 0x25, 0x00), border=GOLD, bw=1)
    txt(s, Inches(0.55), y, Inches(1.3), Inches(0.3),
        label, sz=11, c=color, bold=(i == 0))

# === MAIN CONTENT AREA ===
main_left = Inches(2.3)
main_top = Inches(1.2)

# Header bar
box(s, main_left, main_top, Inches(10.7), Inches(0.55), fill=SIDEBAR_BG, border=RGBColor(0x2a, 0x2a, 0x40))
txt(s, main_left + Inches(0.2), main_top + Inches(0.05), Inches(4), Inches(0.4),
    "สถานีกลางบางซื่อ", sz=14, c=WHITE, bold=True)
txt(s, Inches(10.5), main_top + Inches(0.05), Inches(2.3), Inches(0.4),
    "16 ก.พ. 2569  |  ผู้จัดการก้อง", sz=10, c=DIM, align=PP_ALIGN.RIGHT)

# === KPI CARDS (row 1) ===
kpi_y = Inches(1.95)
kpis = [
    ("2,340", "คนในสถานี", GOLD, "real-time"),
    ("42", "เที่ยวรถไฟวันนี้", BLUE, "จาก SRT"),
    ("3", "แจ้งเหตุรอจัดการ", RED, "urgent"),
    ("18/25", "บูทเปิด/ทั้งหมด", GREEN, "72%"),
    ("1,205", "ผู้ใช้ App วันนี้", PURPLE, "+12%"),
]

for i, (num, label, color, sub) in enumerate(kpis):
    x = main_left + Inches(i * 2.14)
    box(s, x, kpi_y, Inches(2.0), Inches(1.1), fill=CARD_BG, border=color, bw=1.5)
    txt(s, x, kpi_y + Inches(0.1), Inches(2.0), Inches(0.5),
        num, sz=28, c=color, bold=True, align=PP_ALIGN.CENTER)
    txt(s, x, kpi_y + Inches(0.55), Inches(2.0), Inches(0.3),
        label, sz=10, c=LIGHT, align=PP_ALIGN.CENTER)
    txt(s, x, kpi_y + Inches(0.8), Inches(2.0), Inches(0.2),
        sub, sz=8, c=DIM, align=PP_ALIGN.CENTER)

# === HEATMAP (left) ===
heatmap_y = Inches(3.25)
box(s, main_left, heatmap_y, Inches(6.2), Inches(3.1), fill=DARK_CARD, border=GOLD, bw=1)
txt(s, main_left + Inches(0.15), heatmap_y + Inches(0.1), Inches(3), Inches(0.3),
    "Heatmap สถานี (Live)", sz=12, c=GOLD, bold=True)
txt(s, main_left + Inches(5), heatmap_y + Inches(0.1), Inches(1), Inches(0.3),
    "ตอนนี้", sz=9, c=DIM, align=PP_ALIGN.RIGHT)

# Heatmap placeholder zones
zones = [
    (Inches(0.3), Inches(0.6), Inches(1.8), Inches(1.0), RGBColor(0xe7, 0x4c, 0x3c), "ชานชาลา 1", "สูง"),
    (Inches(2.3), Inches(0.6), Inches(1.6), Inches(1.0), RGBColor(0xf3, 0x9c, 0x12), "โถงกลาง", "ปานกลาง"),
    (Inches(4.1), Inches(0.6), Inches(1.8), Inches(1.0), RGBColor(0x27, 0xae, 0x60), "ชานชาลา 2", "ต่ำ"),
    (Inches(0.3), Inches(1.8), Inches(2.5), Inches(0.9), RGBColor(0xf3, 0x9c, 0x12), "ร้านค้าโซน A", "ปานกลาง"),
    (Inches(3.0), Inches(1.8), Inches(2.9), Inches(0.9), RGBColor(0x27, 0xae, 0x60), "ร้านค้าโซน B", "ต่ำ"),
]
for dx, dy, w, h, color, name, level in zones:
    # Zone fill with transparency effect
    zone_fill = RGBColor(
        min(255, color[0] // 3 + BG[0]),
        min(255, color[1] // 3 + BG[1]),
        min(255, color[2] // 3 + BG[2]),
    )
    box(s, main_left + dx, heatmap_y + dy, w, h, fill=zone_fill, border=color, bw=0.8)
    txt(s, main_left + dx, heatmap_y + dy + Inches(0.1), w, Inches(0.3),
        name, sz=9, c=WHITE, bold=True, align=PP_ALIGN.CENTER)
    txt(s, main_left + dx, heatmap_y + dy + h - Inches(0.25), w, Inches(0.2),
        level, sz=8, c=color, align=PP_ALIGN.CENTER)

# === ALERTS (right) ===
alert_x = main_left + Inches(6.45)
box(s, alert_x, heatmap_y, Inches(4.25), Inches(3.1), fill=DARK_CARD, border=RED, bw=1)
txt(s, alert_x + Inches(0.15), heatmap_y + Inches(0.1), Inches(3), Inches(0.3),
    "แจ้งเหตุล่าสุด", sz=12, c=RED, bold=True)

alerts = [
    ("ผู้โดยสารหายของ", "ชานชาลา 3", "5 นาที", RED),
    ("Kiosk #7 ไม่ตอบสนอง", "โถงกลาง", "12 นาที", ORANGE),
    ("ผู้โดยสารเป็นลม", "ทางออก B", "28 นาที", ORANGE),
    ("รถไฟสาย X กลับมาปกติ", "—", "1 ชม.", GREEN),
    ("ไฟโซน C กระพริบ", "ร้านค้า C", "2 ชม.", DIM),
]

for i, (title, loc, time, color) in enumerate(alerts):
    ay = heatmap_y + Inches(0.5 + i * 0.48)
    # Status dot
    dot_color = color
    box(s, alert_x + Inches(0.15), ay + Inches(0.05), Inches(0.15), Inches(0.15),
        fill=dot_color, border=dot_color, bw=0.5)
    txt(s, alert_x + Inches(0.4), ay, Inches(2.2), Inches(0.25),
        title, sz=10, c=WHITE, bold=True)
    txt(s, alert_x + Inches(0.4), ay + Inches(0.22), Inches(1.5), Inches(0.2),
        loc, sz=8, c=DIM)
    txt(s, alert_x + Inches(3.2), ay + Inches(0.05), Inches(0.9), Inches(0.2),
        time, sz=8, c=DIM, align=PP_ALIGN.RIGHT)

# === KIOSK STATUS (bottom) ===
kiosk_y = Inches(6.55)
box(s, main_left, kiosk_y, Inches(10.7), Inches(0.55), fill=DARK_CARD, border=RGBColor(0x2a, 0x2a, 0x40))
txt(s, main_left + Inches(0.15), kiosk_y + Inches(0.05), Inches(2), Inches(0.3),
    "สถานะ Kiosk", sz=10, c=LIGHT, bold=True)

kiosk_statuses = [
    ("#1", GREEN), ("#2", GREEN), ("#3", GREEN), ("#4", GREEN),
    ("#5", ORANGE), ("#6", GREEN), ("#7", RED), ("#8", GREEN),
    ("#9", GREEN), ("#10", GREEN),
]
for i, (name, color) in enumerate(kiosk_statuses):
    kx = main_left + Inches(2.2 + i * 0.85)
    box(s, kx, kiosk_y + Inches(0.1), Inches(0.7), Inches(0.3),
        fill=RGBColor(
            min(255, color[0] // 5 + CARD_BG[0]),
            min(255, color[1] // 5 + CARD_BG[1]),
            min(255, color[2] // 5 + CARD_BG[2]),
        ), border=color, bw=0.8)
    txt(s, kx, kiosk_y + Inches(0.12), Inches(0.7), Inches(0.25),
        name, sz=8, c=color, bold=True, align=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 2 — Design Specs
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
txt(s, Inches(0.5), Inches(0.2), Inches(12), Inches(0.6),
    "Design Specs — สำหรับทีมพัฒนา", sz=28, c=WHITE, bold=True)
txt(s, Inches(0.5), Inches(0.7), Inches(12), Inches(0.4),
    "สี  |  Layout  |  Component  |  หลักคิด", sz=14, c=GOLD)

# Color palette
txt(s, Inches(0.5), Inches(1.3), Inches(3), Inches(0.4),
    "Color Palette", sz=18, c=WHITE, bold=True)

colors_spec = [
    ("#1A1A2E", "Background", BG),
    ("#121222", "Sidebar", SIDEBAR_BG),
    ("#22223A", "Card", CARD_BG),
    ("#E5C100", "Gold (primary)", GOLD),
    ("#E74C3C", "Red (alert)", RED),
    ("#F39C12", "Orange (warning)", ORANGE),
    ("#27AE60", "Green (ok)", GREEN),
    ("#3498DB", "Blue (info)", BLUE),
    ("#9B59B6", "Purple (app)", PURPLE),
]

for i, (hex_val, name, color) in enumerate(colors_spec):
    col = i % 3
    row = i // 3
    cx = Inches(0.5 + col * 2.8)
    cy = Inches(1.8 + row * 0.6)
    box(s, cx, cy, Inches(0.4), Inches(0.4), fill=color, border=LIGHT, bw=0.5)
    txt(s, cx + Inches(0.5), cy + Inches(0.02), Inches(2.2), Inches(0.2),
        hex_val, sz=11, c=WHITE, bold=True)
    txt(s, cx + Inches(0.5), cy + Inches(0.2), Inches(2.2), Inches(0.2),
        name, sz=9, c=DIM)

# Layout rules
txt(s, Inches(0.5), Inches(3.8), Inches(5), Inches(0.4),
    "Layout Rules", sz=18, c=WHITE, bold=True)
bullets(s, Inches(0.5), Inches(4.3), Inches(5), Inches(3), [
    "Sidebar กว้าง ~200px อยู่ซ้ายตลอด",
    "KPI Cards แถวบนสุด — เห็นตัวเลขทันที",
    "Heatmap ครึ่งซ้าย ใหญ่ที่สุด (สำคัญสุด)",
    "แจ้งเหตุ ครึ่งขวา เรียงตาม priority",
    "สถานะ Kiosk แถวล่าง (secondary info)",
    "Font: Tahoma หรือ Inter",
    "ตัวเลข KPI ใหญ่ 28-36px / Label เล็ก 10-12px",
], sz=12, c=LIGHT, sp=6)

# Component list
txt(s, Inches(6.5), Inches(3.8), Inches(6), Inches(0.4),
    "Components ที่ต้องทำ", sz=18, c=WHITE, bold=True)

components = [
    ("KPI Card", "ตัวเลข + label + subtitle\nมีสีขอบตาม type", GOLD),
    ("Heatmap Panel", "Floor plan + gradient overlay\nเปลี่ยน timeframe ได้", GREEN),
    ("Alert List", "เรียง priority สูง→ต่ำ\nกดเพื่อดูรายละเอียด", RED),
    ("Kiosk Status Bar", "แสดงสถานะทุกตู้\nสี = online/warning/offline", BLUE),
    ("Sidebar Menu", "Highlight หน้าปัจจุบัน\nCollapsible", DIM),
]

for i, (name, desc, color) in enumerate(components):
    cy = Inches(4.3 + i * 0.6)
    box(s, Inches(6.5), cy, Inches(0.15), Inches(0.45), fill=color, border=color, bw=0.5)
    txt(s, Inches(6.8), cy, Inches(2), Inches(0.3),
        name, sz=11, c=color, bold=True)
    txt(s, Inches(8.8), cy, Inches(4), Inches(0.45),
        desc.replace('\n', '  |  '), sz=10, c=LIGHT)


# ============================================================
# SLIDE 3 — Role Difference
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
txt(s, Inches(0.5), Inches(0.2), Inches(12), Inches(0.6),
    "2 Roles — หน้าจอเดียวกัน เห็นไม่เหมือนกัน", sz=28, c=WHITE, bold=True)
txt(s, Inches(0.5), Inches(0.7), Inches(12), Inches(0.4),
    "ใช้ Role & Permission คุม — ไม่ต้องทำ 2 หน้าจอ", sz=14, c=GOLD)

# Staff view
box(s, Inches(0.3), Inches(1.3), Inches(6.2), Inches(5.8), fill=DARK_CARD, border=BLUE, bw=2)
txt(s, Inches(0.3), Inches(1.4), Inches(6.2), Inches(0.5),
    "เจ้าหน้าที่ เห็น:", sz=20, c=BLUE, bold=True, align=PP_ALIGN.CENTER)

staff_sees = [
    ("KPI Cards", "คนในสถานี, แจ้งเหตุ, สถานะ Kiosk", GREEN),
    ("Heatmap สถานี", "ดูความหนาแน่น real-time", GREEN),
    ("แจ้งเหตุล่าสุด", "รับเรื่อง + จัดการ (งานหลัก)", GREEN),
    ("สถานะ Kiosk", "ตู้ไหน offline ต้องไปดู", GREEN),
]
staff_not = [
    ("Revenue / รายได้", RED),
    ("Executive Report", RED),
    ("Export ข้อมูล", RED),
    ("User Management", RED),
    ("Role & Permission", RED),
]

for i, (item, desc, color) in enumerate(staff_sees):
    y = Inches(2.1 + i * 0.55)
    txt(s, Inches(0.5), y, Inches(2.2), Inches(0.3),
        f"  {item}", sz=12, c=GREEN, bold=True)
    txt(s, Inches(2.7), y, Inches(3.5), Inches(0.3),
        desc, sz=11, c=LIGHT)

txt(s, Inches(0.5), Inches(4.4), Inches(5.8), Inches(0.3),
    "เข้าไม่ได้:", sz=12, c=RED, bold=True)
for i, (item, color) in enumerate(staff_not):
    y = Inches(4.8 + i * 0.4)
    txt(s, Inches(0.5), y, Inches(5.8), Inches(0.3),
        f"  {item}", sz=11, c=RED)

# Admin view
box(s, Inches(6.8), Inches(1.3), Inches(6.2), Inches(5.8), fill=DARK_CARD, border=GOLD, bw=2)
txt(s, Inches(6.8), Inches(1.4), Inches(6.2), Inches(0.5),
    "ผู้บริหาร เห็น:", sz=20, c=GOLD, bold=True, align=PP_ALIGN.CENTER)

admin_sees = [
    ("ทุกอย่างที่เจ้าหน้าที่เห็น", "+", GREEN),
    ("Revenue Dashboard", "รายได้รวม, แยก Tier/โซน", GOLD),
    ("Executive Report", "สรุปสำหรับผู้บริหาร", GOLD),
    ("Export (PDF/Excel)", "ส่ง report อัตโนมัติ", GOLD),
    ("User Management", "จัดการ user ทั้งระบบ", GOLD),
    ("Role & Permission", "กำหนดสิทธิ์ทุก role", GOLD),
    ("Kiosk Config", "ตั้งค่า/อัพเดท Kiosk ทุกตู้", GOLD),
    ("Analytics ขั้นสูง", "Trend, forecast, compare", GOLD),
]

for i, (item, desc, color) in enumerate(admin_sees):
    y = Inches(2.1 + i * 0.55)
    txt(s, Inches(7.0), y, Inches(2.5), Inches(0.3),
        f"  {item}", sz=12, c=color, bold=True)
    txt(s, Inches(9.5), y, Inches(3.2), Inches(0.3),
        desc, sz=11, c=LIGHT)

txt(s, Inches(6.8), Inches(6.5), Inches(6.2), Inches(0.4),
    "เข้าถึงทุกเมนู ทุกฟังก์ชัน", sz=12, c=GOLD, bold=True, align=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 4 — Sidebar Menu Detail
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
txt(s, Inches(0.5), Inches(0.2), Inches(12), Inches(0.6),
    "เมนูทั้งหมดใน Core Dashboard", sz=28, c=WHITE, bold=True)
txt(s, Inches(0.5), Inches(0.7), Inches(12), Inches(0.4),
    "Sidebar เมนู → แต่ละหน้าทำอะไรบ้าง", sz=14, c=GOLD)

pages = [
    ("Overview", "หน้าแรก", [
        "KPI Cards — ตัวเลขสำคัญ",
        "Heatmap สถานี (Live)",
        "แจ้งเหตุล่าสุด",
        "สถานะ Kiosk ทุกตู้",
    ], GOLD),
    ("Users", "จัดการผู้ใช้", [
        "รายชื่อ User ทั้งระบบ",
        "เพิ่ม / ลบ / แก้ไข",
        "Ban / Suspend",
        "ดูประวัติการใช้งาน",
    ], BLUE),
    ("Roles", "จัดการสิทธิ์", [
        "สร้าง/แก้ไข Role",
        "กำหนด Permission",
        "ดู Audit Log",
    ], PURPLE),
    ("Report", "รายงาน", [
        "Executive Report",
        "Traffic Analytics",
        "Revenue Dashboard",
        "Export PDF / Excel",
    ], GREEN),
    ("Alert", "แจ้งเหตุ", [
        "รายการแจ้งเหตุทั้งหมด",
        "Filter ตาม status/priority",
        "ตอบรับ / ปิดเคส",
        "ประวัติแจ้งเหตุ",
    ], RED),
    ("Kiosk", "จัดการ Kiosk", [
        "สถานะ Kiosk ทุกตู้",
        "ตั้งค่า / อัพเดทเนื้อหา",
        "จัดการโฆษณา / Popup",
        "ดู Analytics การใช้งาน",
    ], ORANGE),
]

for i, (name, thai_name, items, color) in enumerate(pages):
    col = i % 3
    row = i // 3
    x = Inches(0.3 + col * 4.3)
    y = Inches(1.2 + row * 3.1)
    box(s, x, y, Inches(4.1), Inches(2.8), fill=DARK_CARD, border=color, bw=2)
    txt(s, x + Inches(0.15), y + Inches(0.1), Inches(3.8), Inches(0.4),
        f"{name} — {thai_name}", sz=16, c=color, bold=True)
    bullets(s, x + Inches(0.15), y + Inches(0.6), Inches(3.8), Inches(2),
            items, sz=12, c=LIGHT, sp=6)


# ============================================================
# SAVE
# ============================================================
prs.save(OUTPUT_PATH)
print(f"Saved: {OUTPUT_PATH}")
print(f"Slides: {len(prs.slides)}")
