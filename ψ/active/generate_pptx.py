"""
Generate SWT Presentation .pptx (v3 — NDF scope, demo-ready)
NDF: Core Dashboard + Core DB + Auth + Kiosk + Mobile App + Admin Dashboard
Other teams: Space Management, IoT People Counting
Client data: SRT API (การรถไฟ)
"""

import sys
sys.stdout.reconfigure(encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
OUTPUT_PATH = os.path.join(SCRIPT_DIR, "SWT-CoreDashboard-Kiosk.pptx")

# Brand colors
DARK_BG = RGBColor(0x1a, 0x1a, 0x2e)
ACCENT_GOLD = RGBColor(0xe5, 0xc1, 0x00)
ACCENT_ORANGE = RGBColor(0xf3, 0x9c, 0x12)
ACCENT_BLUE = RGBColor(0x34, 0x98, 0xdb)
ACCENT_GREEN = RGBColor(0x27, 0xae, 0x60)
ACCENT_RED = RGBColor(0xe7, 0x4c, 0x3c)
ACCENT_PURPLE = RGBColor(0x9b, 0x59, 0xb6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xBD, 0xC3, 0xC7)
DARK_GRAY = RGBColor(0x2c, 0x3e, 0x50)
DIM_GRAY = RGBColor(0x7f, 0x8c, 0x8d)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def add_bg(slide, color=DARK_BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color=None, border_color=None, border_width=1.5):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color or RGBColor(0x22, 0x22, 0x3a)
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(border_width)
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, left, top, width, height, text, size=18, color=WHITE,
             bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = 'Tahoma'
    p.alignment = align
    return tb


def add_bullets(slide, left, top, width, height, items, size=15, color=LIGHT_GRAY, spacing=8):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = 'Tahoma'
        p.space_after = Pt(spacing)
    return tb


def add_screenshot_placeholder(slide, left, top, width, height, label, border_color=LIGHT_GRAY):
    """Add a mockup browser/screen frame with placeholder text"""
    # Top bar
    add_shape(slide, left, top, width, Inches(0.3),
              fill_color=RGBColor(0x33, 0x33, 0x44), border_color=border_color)
    add_text(slide, left + Inches(0.2), top + Inches(0.02), Inches(4), Inches(0.3),
             f"●  ●  ●   {label}", size=10, color=DIM_GRAY)
    # Screen area
    add_shape(slide, left, top + Inches(0.3), width, height - Inches(0.3),
              fill_color=RGBColor(0x12, 0x12, 0x20), border_color=border_color)
    mid_y = top + Inches(0.3) + (height - Inches(0.3)) / 2 - Inches(0.3)
    add_text(slide, left, mid_y, width, Inches(0.6),
             f"[ วางภาพแคป {label} ตรงนี้ ]",
             size=20, color=DIM_GRAY, align=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 1 — ปก
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_shape(s, Inches(1), Inches(2.8), Inches(3), Pt(4), fill_color=ACCENT_GOLD)
add_text(s, Inches(1), Inches(1.2), Inches(11), Inches(1.2),
         "Smart Tourism Web Application", size=44, color=WHITE, bold=True)
add_text(s, Inches(1), Inches(2.1), Inches(11), Inches(0.8),
         "Core System — by NDF", size=28, color=ACCENT_GOLD)
add_text(s, Inches(1), Inches(3.2), Inches(11), Inches(0.6),
         "Core Dashboard  |  Core Database  |  Kiosk  |  Mobile App  |  Admin Dashboard", size=18, color=LIGHT_GRAY)
add_text(s, Inches(1), Inches(4.5), Inches(11), Inches(0.5),
         "ระบบจัดการสถานีอัจฉริยะ — ครบวงจร ทุกฟังก์ชัน ทุก User Group", size=16, color=DIM_GRAY)
add_text(s, Inches(1), Inches(6.2), Inches(11), Inches(0.5),
         "NDF Co., Ltd.", size=18, color=ACCENT_GOLD, bold=True)


# ============================================================
# SLIDE 2 — ภาพรวมระบบของ NDF
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_text(s, Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
         "ภาพรวมระบบ — NDF รับผิดชอบอะไร", size=32, color=WHITE, bold=True)
add_text(s, Inches(0.5), Inches(0.9), Inches(12), Inches(0.5),
         "Core System เป็นหัวใจ — เชื่อมทุกระบบเข้าด้วยกัน", size=18, color=ACCENT_GOLD)

# NDF scope (large box)
add_shape(s, Inches(0.3), Inches(1.6), Inches(8.5), Inches(5.4),
          fill_color=RGBColor(0x1e, 0x20, 0x30), border_color=ACCENT_GOLD, border_width=3)
add_text(s, Inches(0.3), Inches(1.3), Inches(8.5), Inches(0.4),
         "NDF — ขอบเขตงานของเรา", size=14, color=ACCENT_GOLD, bold=True, align=PP_ALIGN.CENTER)

# Core block
add_shape(s, Inches(2.5), Inches(1.9), Inches(4.2), Inches(2.2),
          fill_color=RGBColor(0x2a, 0x25, 0x00), border_color=ACCENT_GOLD, border_width=2)
add_text(s, Inches(2.5), Inches(2.0), Inches(4.2), Inches(0.5),
         "Core Dashboard + Database", size=18, color=ACCENT_GOLD, bold=True, align=PP_ALIGN.CENTER)
add_bullets(s, Inches(2.8), Inches(2.5), Inches(3.6), Inches(1.5), [
    "Auth กลาง — ออก Token ให้ทุกระบบ",
    "Report & Analytics",
    "User & Role Management",
    "Core Database — ศูนย์ข้อมูลรวม",
], size=12, color=LIGHT_GRAY)

# Frontend blocks
frontends = [
    ("ตู้ Kiosk", Inches(0.5), Inches(4.5), ACCENT_ORANGE),
    ("Mobile App", Inches(3.1), Inches(4.5), ACCENT_PURPLE),
    ("Admin Dashboard", Inches(5.7), Inches(4.5), ACCENT_BLUE),
]
for name, x, y, color in frontends:
    add_shape(s, x, y, Inches(2.4), Inches(1.0), border_color=color)
    add_text(s, x, y + Inches(0.05), Inches(2.4), Inches(0.4),
             name, size=14, color=color, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, x, y + Inches(0.5), Inches(2.4), Inches(0.4),
             "ทีมเรา", size=11, color=DIM_GRAY, align=PP_ALIGN.CENTER)

# API blocks
apis = [
    ("Kiosk API", Inches(0.5), Inches(5.7), ACCENT_ORANGE),
    ("Web/App API", Inches(3.1), Inches(5.7), ACCENT_PURPLE),
    ("LiveWalk API", Inches(5.7), Inches(5.7), ACCENT_GREEN),
]
for name, x, y, color in apis:
    add_shape(s, x, y, Inches(2.4), Inches(0.7), border_color=color)
    add_text(s, x, y + Inches(0.05), Inches(2.4), Inches(0.4),
             name, size=12, color=color, bold=True, align=PP_ALIGN.CENTER)

# External systems (right side)
add_text(s, Inches(9.2), Inches(1.3), Inches(4), Inches(0.4),
         "ระบบภายนอก", size=14, color=DIM_GRAY, bold=True, align=PP_ALIGN.CENTER)

externals = [
    ("Space System", "ทีมอื่น", Inches(9.2), Inches(1.8), DIM_GRAY),
    ("IoT API", "ทีมอื่น", Inches(9.2), Inches(3.0), DIM_GRAY),
    ("SRT API", "การรถไฟ (ลูกค้า)", Inches(9.2), Inches(4.2), DIM_GRAY),
]
for name, team, x, y, color in externals:
    add_shape(s, x, y, Inches(3.8), Inches(1.0), border_color=color)
    add_text(s, x, y + Inches(0.05), Inches(3.8), Inches(0.4),
             name, size=14, color=color, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, x, y + Inches(0.5), Inches(3.8), Inches(0.4),
             f"({team})", size=11, color=DIM_GRAY, align=PP_ALIGN.CENTER)

# Arrows
add_text(s, Inches(8.5), Inches(2.5), Inches(0.8), Inches(0.5),
         "←→", size=20, color=ACCENT_GOLD, bold=True, align=PP_ALIGN.CENTER)
add_text(s, Inches(8.5), Inches(3.5), Inches(0.8), Inches(0.5),
         "←→", size=20, color=ACCENT_GOLD, bold=True, align=PP_ALIGN.CENTER)
add_text(s, Inches(8.5), Inches(4.5), Inches(0.8), Inches(0.5),
         "←→", size=20, color=ACCENT_GOLD, bold=True, align=PP_ALIGN.CENTER)

# Integration API label
add_shape(s, Inches(9.2), Inches(5.5), Inches(3.8), Inches(0.7),
          fill_color=RGBColor(0x2a, 0x25, 0x00), border_color=ACCENT_GOLD)
add_text(s, Inches(9.2), Inches(5.55), Inches(3.8), Inches(0.5),
         "Integration API — เชื่อมทุกระบบ", size=12, color=ACCENT_GOLD, align=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 3 — User Groups
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_text(s, Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
         "User Groups — ใครใช้อะไร", size=32, color=WHITE, bold=True)
add_text(s, Inches(0.5), Inches(0.9), Inches(12), Inches(0.5),
         "4 กลุ่มผู้ใช้ — NDF ดูแล 3 กลุ่ม ผู้เช่าอยู่ฝั่ง Space System", size=18, color=ACCENT_GOLD)

users = [
    ("นักท่องเที่ยว / ผู้โดยสาร", [
        "ใช้ตู้ Kiosk ที่สถานี",
        "ใช้ Mobile App (iOS/Android)",
        "ดูตารางรถไฟ จองตั๋ว",
        "นำทาง ค้นหาร้านค้า",
        "สะสม Point / Reward",
    ], ACCENT_ORANGE, "Kiosk + Mobile App"),
    ("เจ้าหน้าที่", [
        "เข้า Core Dashboard (บางส่วน)",
        "รับแจ้งเหตุ / ช่วยเหลือ",
        "ดู Report เบื้องต้น",
        "จัดการ Kiosk",
    ], ACCENT_BLUE, "Core Dashboard"),
    ("ผู้บริหารระบบ", [
        "เข้า Core Dashboard (ทั้งหมด)",
        "จัดการข้อมูลร้านค้า ผังสถานี",
        "กำหนดสิทธิ์ Role & Permission",
        "ดู Executive Report",
        "Export ข้อมูล",
    ], ACCENT_GOLD, "Core Dashboard"),
    ("ผู้เช่า / ร้านค้า", [
        "ลงทะเบียน จองพื้นที่",
        "ชำระเงิน ทำสัญญา",
        "ดูสถิติหน้าร้าน",
        "(ใช้ Auth จาก Core ของเรา)",
    ], DIM_GRAY, "Space System (ทีมอื่น)"),
]

for i, (role, items, color, access) in enumerate(users):
    x = Inches(0.3 + i * 3.2)
    add_shape(s, x, Inches(1.6), Inches(3), Inches(5.2), border_color=color, border_width=2)
    add_text(s, x, Inches(1.7), Inches(3), Inches(0.5),
             role, size=16, color=color, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, x, Inches(2.2), Inches(3), Inches(0.4),
             f"เข้าผ่าน: {access}", size=11, color=DIM_GRAY, align=PP_ALIGN.CENTER)
    add_bullets(s, x + Inches(0.2), Inches(2.7), Inches(2.6), Inches(3.5),
                items, size=12, color=LIGHT_GRAY, spacing=8)


# ============================================================
# SLIDE 4 — ตู้ Kiosk: ภาพรวม
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_text(s, Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
         "ตู้ Kiosk — ภาพรวมฟังก์ชัน", size=32, color=WHITE, bold=True)
add_text(s, Inches(0.5), Inches(0.9), Inches(12), Inches(0.5),
         "จุดให้บริการข้อมูลอัจฉริยะ ติดตั้งในสถานี — สำหรับนักท่องเที่ยว/ผู้โดยสาร", size=18, color=ACCENT_ORANGE)

kiosk_features = [
    ("AI Character", "ตัวละคร AI ต้อนรับ\nดูผ่านหน้าจอ/Smartphone\nโต้ตอบภาษาธรรมชาติ", ACCENT_ORANGE),
    ("Voice แนะนำ", "STT + TTS หลายภาษา\nไทย/อังกฤษ/จีน/ญี่ปุ่น\nถาม-ตอบด้วยเสียง", ACCENT_GOLD),
    ("Tourism / Event Info", "ข้อมูลท่องเที่ยวในพื้นที่\nกิจกรรม อีเวนต์\nอัพเดทจาก Admin", ACCENT_BLUE),
    ("Route Finder", "นำทางภายในสถานี\nจากจุดปัจจุบัน → ปลายทาง\nInteractive map", ACCENT_GREEN),
    ("Emergency Call", "โทรฉุกเฉิน / แจ้งเหตุ\nเชื่อมหน่วยงานโดยตรง\nส่ง alert ถึง Staff", ACCENT_RED),
    ("Train Schedule", "ตารางรถไฟ real-time\nข้อมูลจาก SRT API\nเวลา / สาย / สถานี", ACCENT_PURPLE),
]

for i, (title, desc, color) in enumerate(kiosk_features):
    col = i % 3
    row = i // 3
    x = Inches(0.5 + col * 4.15)
    y = Inches(1.6 + row * 2.8)
    add_shape(s, x, y, Inches(3.9), Inches(2.4), border_color=color)
    add_text(s, x + Inches(0.2), y + Inches(0.15), Inches(3.5), Inches(0.5),
             title, size=18, color=color, bold=True)
    add_text(s, x + Inches(0.2), y + Inches(0.7), Inches(3.5), Inches(1.5),
             desc, size=13, color=LIGHT_GRAY)


# ============================================================
# SLIDE 5 — ตู้ Kiosk: AI Assistant Demo
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_text(s, Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
         "ตู้ Kiosk — AI Assistant Demo", size=32, color=WHITE, bold=True)
add_text(s, Inches(0.5), Inches(0.9), Inches(12), Inches(0.5),
         "สาธิตการถาม-ตอบผ่าน AI Character + Voice", size=18, color=ACCENT_ORANGE)

add_screenshot_placeholder(s, Inches(0.5), Inches(1.5), Inches(6), Inches(5.2),
                           "Kiosk — AI Conversation", ACCENT_ORANGE)

# Right side: feature callouts
add_text(s, Inches(7), Inches(1.8), Inches(5.5), Inches(0.5),
         "สิ่งที่เห็นในหน้าจอ:", size=18, color=ACCENT_ORANGE, bold=True)
add_bullets(s, Inches(7), Inches(2.4), Inches(5.5), Inches(4), [
    "AI Character แสดงบนหน้าจอ",
    "ผู้ใช้พิมพ์หรือพูดคำถาม",
    "AI ตอบกลับเป็นข้อความ + เสียง",
    "แนะนำเส้นทาง ร้านอาหาร สถานที่",
    "รองรับหลายภาษา",
    "",
    "Kiosk API Backend:",
    "• AI/LLM Engine ประมวลผลคำถาม",
    "• Voice STT/TTS แปลงเสียง",
    "• FAQ/Tourism Info ฐานข้อมูล",
], size=13, color=LIGHT_GRAY, spacing=8)


# ============================================================
# SLIDE 6 — ตู้ Kiosk: Emergency & Route Finder
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_text(s, Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
         "ตู้ Kiosk — Emergency & Route Finder", size=32, color=WHITE, bold=True)
add_text(s, Inches(0.5), Inches(0.9), Inches(12), Inches(0.5),
         "แจ้งเหตุฉุกเฉิน + นำทางในสถานี", size=18, color=ACCENT_RED)

# Left: Emergency
add_screenshot_placeholder(s, Inches(0.5), Inches(1.5), Inches(5.8), Inches(5.2),
                           "Kiosk — Emergency Call", ACCENT_RED)

# Right: Route Finder
add_screenshot_placeholder(s, Inches(6.8), Inches(1.5), Inches(5.8), Inches(5.2),
                           "Kiosk — Route Finder", ACCENT_GREEN)


# ============================================================
# SLIDE 7 — Mobile App: ภาพรวม
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_text(s, Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
         "Mobile App — ภาพรวมฟังก์ชัน", size=32, color=WHITE, bold=True)
add_text(s, Inches(0.5), Inches(0.9), Inches(12), Inches(0.5),
         "iOS & Android — สำหรับนักท่องเที่ยว/ผู้โดยสาร ใช้ได้ทุกที่", size=18, color=ACCENT_PURPLE)

mobile_features = [
    ("ThaiD / Login", "ยืนยันตัวตนผ่าน ThaiD\nEmail, Social Media", ACCENT_PURPLE),
    ("Train Schedule & Booking", "ตารางรถไฟ real-time\nจองตั๋ว ชำระเงิน", ACCENT_BLUE),
    ("Navigation", "นำทาง indoor/outdoor\nแผนที่สถานี interactive", ACCENT_GREEN),
    ("Mini-ChatBot", "AI ตอบคำถามเบื้องต้น\nช่วยหาข้อมูล", ACCENT_ORANGE),
    ("Notification", "แจ้งเตือน Push / In-app\nรถไฟล่าช้า โปรโมชั่น", ACCENT_RED),
    ("Point / Reward", "สะสมแต้ม แลกรางวัล\nHistory การใช้แต้ม", ACCENT_GOLD),
    ("Favorites", "บันทึกสถานที่/ร้านค้าชอบ\nเข้าถึงเร็ว", ACCENT_BLUE),
    ("Search", "ค้นหาร้านค้า สถานที่\nค้นหาต่อเนื่อง", ACCENT_PURPLE),
    ("Emergency", "แจ้งเหตุร้าย\nส่ง alert ถึงเจ้าหน้าที่", ACCENT_RED),
    ("Payment / Refund", "ชำระเงินในแอป\nขอคืนเงิน", ACCENT_GREEN),
    ("History", "ประวัติการจอง\nประวัติการแจ้งเหตุ", ACCENT_ORANGE),
    ("Accessibility", "รองรับหลายภาษา\nThaiD verification", DIM_GRAY),
]

for i, (title, desc, color) in enumerate(mobile_features):
    col = i % 4
    row = i // 4
    x = Inches(0.3 + col * 3.2)
    y = Inches(1.6 + row * 1.9)
    add_shape(s, x, y, Inches(3), Inches(1.6), border_color=color)
    add_text(s, x + Inches(0.15), y + Inches(0.1), Inches(2.7), Inches(0.4),
             title, size=14, color=color, bold=True)
    add_text(s, x + Inches(0.15), y + Inches(0.55), Inches(2.7), Inches(0.9),
             desc, size=12, color=LIGHT_GRAY)


# ============================================================
# SLIDE 8 — Mobile App: Train Schedule & Booking Demo
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_text(s, Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
         "Mobile App — Train Schedule & Booking", size=32, color=WHITE, bold=True)
add_text(s, Inches(0.5), Inches(0.9), Inches(12), Inches(0.5),
         "ดูตารางรถไฟ real-time + จองตั๋ว (ข้อมูลจาก SRT API)", size=18, color=ACCENT_PURPLE)

add_screenshot_placeholder(s, Inches(0.5), Inches(1.5), Inches(5.8), Inches(5.2),
                           "App — Train Schedule", ACCENT_PURPLE)
add_screenshot_placeholder(s, Inches(6.8), Inches(1.5), Inches(5.8), Inches(5.2),
                           "App — Booking Flow", ACCENT_BLUE)


# ============================================================
# SLIDE 9 — Mobile App: Navigation & ChatBot Demo
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_text(s, Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
         "Mobile App — Navigation & ChatBot", size=32, color=WHITE, bold=True)
add_text(s, Inches(0.5), Inches(0.9), Inches(12), Inches(0.5),
         "นำทางในสถานี + Mini-ChatBot ถาม-ตอบ", size=18, color=ACCENT_GREEN)

add_screenshot_placeholder(s, Inches(0.5), Inches(1.5), Inches(5.8), Inches(5.2),
                           "App — Navigation", ACCENT_GREEN)
add_screenshot_placeholder(s, Inches(6.8), Inches(1.5), Inches(5.8), Inches(5.2),
                           "App — ChatBot", ACCENT_ORANGE)


# ============================================================
# SLIDE 10 — Mobile App: Point/Reward & Favorites Demo
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_text(s, Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
         "Mobile App — Point/Reward & Favorites", size=32, color=WHITE, bold=True)
add_text(s, Inches(0.5), Inches(0.9), Inches(12), Inches(0.5),
         "สะสมแต้ม แลกของรางวัล + บันทึกสถานที่ชอบ", size=18, color=ACCENT_GOLD)

add_screenshot_placeholder(s, Inches(0.5), Inches(1.5), Inches(5.8), Inches(5.2),
                           "App — Point & Reward", ACCENT_GOLD)
add_screenshot_placeholder(s, Inches(6.8), Inches(1.5), Inches(5.8), Inches(5.2),
                           "App — Favorites", ACCENT_BLUE)


# ============================================================
# SLIDE 11 — Mobile App: Notification & Emergency Demo
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_text(s, Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
         "Mobile App — Notification & Emergency", size=32, color=WHITE, bold=True)
add_text(s, Inches(0.5), Inches(0.9), Inches(12), Inches(0.5),
         "ระบบแจ้งเตือน + แจ้งเหตุร้าย", size=18, color=ACCENT_RED)

add_screenshot_placeholder(s, Inches(0.5), Inches(1.5), Inches(5.8), Inches(5.2),
                           "App — Notification", ACCENT_ORANGE)
add_screenshot_placeholder(s, Inches(6.8), Inches(1.5), Inches(5.8), Inches(5.2),
                           "App — Emergency Report", ACCENT_RED)


# ============================================================
# SLIDE 12 — Core Dashboard: หน้าแรก
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_text(s, Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
         "Core Dashboard — หน้าแรก (Overview)", size=32, color=WHITE, bold=True)
add_text(s, Inches(0.5), Inches(0.9), Inches(12), Inches(0.5),
         "เปิดมาเห็นสถานการณ์สถานีทันที — สำหรับเจ้าหน้าที่ + ผู้บริหาร", size=18, color=ACCENT_GOLD)

# Dashboard overview mockup
add_screenshot_placeholder(s, Inches(0.5), Inches(1.5), Inches(8), Inches(5.2),
                           "Core Dashboard — Overview", ACCENT_GOLD)

# Right: Key metrics explained
add_text(s, Inches(8.8), Inches(1.8), Inches(4), Inches(0.5),
         "ข้อมูลที่แสดง:", size=18, color=ACCENT_GOLD, bold=True)
add_bullets(s, Inches(8.8), Inches(2.4), Inches(4), Inches(4.5), [
    "จำนวนคนในสถานี (real-time)",
    "จำนวนเที่ยวรถไฟวันนี้",
    "แจ้งเหตุที่รอจัดการ",
    "สถานะ Kiosk ทุกตู้",
    "ผู้ใช้ App วันนี้",
    "บูทที่เปิดอยู่",
    "",
    "Heatmap สถานี (Live)",
    "แจ้งเหตุล่าสุด",
], size=13, color=LIGHT_GRAY, spacing=8)

# Role difference note
add_shape(s, Inches(0.5), Inches(6.8), Inches(12.3), Inches(0.5),
          fill_color=RGBColor(0x1e, 0x1e, 0x30), border_color=ACCENT_GOLD)
add_text(s, Inches(0.5), Inches(6.85), Inches(12.3), Inches(0.4),
         "เจ้าหน้าที่เห็น: แจ้งเหตุ + Heatmap + สถานะ Kiosk  |  ผู้บริหารเห็น: ทุกอย่าง + Revenue + Export",
         size=13, color=ACCENT_GOLD, align=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 13 — Core Dashboard: User & Role Management
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_text(s, Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
         "Core Dashboard — User & Role Management", size=32, color=WHITE, bold=True)
add_text(s, Inches(0.5), Inches(0.9), Inches(12), Inches(0.5),
         "จัดการผู้ใช้ + กำหนดสิทธิ์ตาม Role", size=18, color=ACCENT_BLUE)

# Left: User Management
add_shape(s, Inches(0.5), Inches(1.6), Inches(5.8), Inches(5.2),
          border_color=ACCENT_BLUE, border_width=2)
add_text(s, Inches(0.5), Inches(1.7), Inches(5.8), Inches(0.5),
         "User Management", size=22, color=ACCENT_BLUE, bold=True, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.5), Inches(2.2), Inches(5.8), Inches(0.4),
         "\"ใครอยู่ในระบบบ้าง?\"", size=14, color=DIM_GRAY, align=PP_ALIGN.CENTER)
add_bullets(s, Inches(0.8), Inches(2.8), Inches(5.2), Inches(3.5), [
    "ดูรายชื่อ User ทั้งระบบ",
    "เพิ่ม / ลบ / แก้ไขข้อมูล User",
    "Ban / Suspend User",
    "ดูประวัติการใช้งาน",
    "Migrate User จากแอปเดิม (SRT)",
], size=14, color=LIGHT_GRAY, spacing=10)

add_screenshot_placeholder(s, Inches(0.8), Inches(5.0), Inches(5.2), Inches(1.6),
                           "User Management", ACCENT_BLUE)

# Right: Role & Permission
add_shape(s, Inches(6.8), Inches(1.6), Inches(5.8), Inches(5.2),
          border_color=ACCENT_GOLD, border_width=2)
add_text(s, Inches(6.8), Inches(1.7), Inches(5.8), Inches(0.5),
         "Role & Permission", size=22, color=ACCENT_GOLD, bold=True, align=PP_ALIGN.CENTER)
add_text(s, Inches(6.8), Inches(2.2), Inches(5.8), Inches(0.4),
         "\"คนนี้ทำอะไรได้บ้าง?\"", size=14, color=DIM_GRAY, align=PP_ALIGN.CENTER)
add_bullets(s, Inches(7.1), Inches(2.8), Inches(5.2), Inches(3.5), [
    "สร้าง Role: ผู้บริหาร, เจ้าหน้าที่, Staff",
    "กำหนด Permission แต่ละ Role",
    "เจ้าหน้าที่ดู Report ได้ แต่ลบ User ไม่ได้",
    "Audit Log — ใครทำอะไร เมื่อไหร่",
], size=14, color=LIGHT_GRAY, spacing=10)

add_screenshot_placeholder(s, Inches(7.1), Inches(5.0), Inches(5.2), Inches(1.6),
                           "Role & Permission", ACCENT_GOLD)


# ============================================================
# SLIDE 14 — Core Dashboard: Report & Analytics
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_text(s, Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
         "Core Dashboard — Report & Analytics", size=32, color=WHITE, bold=True)
add_text(s, Inches(0.5), Inches(0.9), Inches(12), Inches(0.5),
         "รายงานสรุปจากทุกระบบ — ข้อมูล IoT + Space + App รวมที่เดียว", size=18, color=ACCENT_GOLD)

reports = [
    ("Executive Report", "สรุปสำหรับผู้บริหาร\nตัวเลขสำคัญ + แนวโน้ม\nเปรียบเทียบ period", ACCENT_GOLD),
    ("Traffic Analytics", "ข้อมูลจาก IoT sensors\nHeatmap, Peak Hours\nจำนวนคนต่อโซน", ACCENT_ORANGE),
    ("Revenue Dashboard", "รายได้รวมทุกบูท\nอัตราการเติบโต\nRevenue forecast", ACCENT_GREEN),
    ("Export", "PDF / Excel\nEmail อัตโนมัติ\nตั้งเวลาส่ง report", ACCENT_BLUE),
]

for i, (title, desc, color) in enumerate(reports):
    x = Inches(0.3 + i * 3.2)
    add_shape(s, x, Inches(1.6), Inches(3), Inches(2.3), border_color=color)
    add_text(s, x + Inches(0.2), Inches(1.7), Inches(2.6), Inches(0.5),
             title, size=18, color=color, bold=True)
    add_text(s, x + Inches(0.2), Inches(2.3), Inches(2.6), Inches(1.2),
             desc, size=13, color=LIGHT_GRAY)

add_screenshot_placeholder(s, Inches(0.5), Inches(4.2), Inches(12.3), Inches(2.8),
                           "Report & Analytics Dashboard", ACCENT_GOLD)


# ============================================================
# SLIDE 15 — Core Dashboard: Auth กลาง
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_text(s, Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
         "Auth กลาง — ระบบยืนยันตัวตนเดียวใช้ทุกระบบ", size=30, color=WHITE, bold=True)
add_text(s, Inches(0.5), Inches(0.9), Inches(12), Inches(0.5),
         "สมัครทีเดียว Login ได้ทุกที่ — ปลอดภัย ควบคุมจากจุดเดียว", size=18, color=ACCENT_RED)

# Center: Auth API
add_shape(s, Inches(4.5), Inches(1.8), Inches(4.3), Inches(2.8),
          fill_color=RGBColor(0x3d, 0x15, 0x15), border_color=ACCENT_RED, border_width=3)
add_text(s, Inches(4.5), Inches(1.9), Inches(4.3), Inches(0.5),
         "Core Auth API", size=22, color=ACCENT_RED, bold=True, align=PP_ALIGN.CENTER)
add_bullets(s, Inches(4.8), Inches(2.5), Inches(3.7), Inches(1.8), [
    "Login: Email / Social / ThaiD",
    "JWT Token Management",
    "Role-based Access Control",
    "Session & Audit Log",
], size=14, color=LIGHT_GRAY)

# Connected systems
systems = [
    ("ตู้ Kiosk", "NDF", Inches(0.5), Inches(2.0), ACCENT_ORANGE),
    ("Mobile App", "NDF", Inches(0.5), Inches(3.3), ACCENT_PURPLE),
    ("Admin Dashboard", "NDF", Inches(0.5), Inches(4.6), ACCENT_BLUE),
    ("Space System", "ทีมอื่น", Inches(9.5), Inches(2.0), DIM_GRAY),
    ("IoT API", "ทีมอื่น", Inches(9.5), Inches(3.3), DIM_GRAY),
    ("SRT API", "การรถไฟ", Inches(9.5), Inches(4.6), DIM_GRAY),
]

for name, team, x, y, color in systems:
    add_shape(s, x, y, Inches(3.5), Inches(1.0), border_color=color)
    add_text(s, x, y + Inches(0.05), Inches(3.5), Inches(0.4),
             name, size=15, color=color, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, x, y + Inches(0.5), Inches(3.5), Inches(0.3),
             f"({team})", size=11, color=DIM_GRAY, align=PP_ALIGN.CENTER)

# Arrows
for y in [Inches(2.3), Inches(3.6), Inches(4.9)]:
    add_text(s, Inches(3.8), y, Inches(0.8), Inches(0.4),
             "→", size=24, color=ACCENT_RED, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, Inches(8.7), y, Inches(0.8), Inches(0.4),
             "←", size=24, color=ACCENT_RED, bold=True, align=PP_ALIGN.CENTER)

# Login flow
add_shape(s, Inches(0.5), Inches(6.0), Inches(12.3), Inches(1.2),
          fill_color=RGBColor(0x1e, 0x1e, 0x30), border_color=ACCENT_RED)
add_text(s, Inches(0.8), Inches(6.1), Inches(11), Inches(0.4),
         "ตัวอย่าง Flow:", size=14, color=ACCENT_RED, bold=True)
add_text(s, Inches(0.8), Inches(6.5), Inches(11), Inches(0.5),
         "สมชายสมัครที่ Mobile App → Auth API บันทึกใน Core DB → วันต่อมาแตะ Kiosk → Auth API ยืนยัน → ใช้ได้เลย ไม่ต้องสมัครใหม่",
         size=14, color=LIGHT_GRAY)


# ============================================================
# SLIDE 16 — การเชื่อมต่อระบบภายนอก
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_text(s, Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
         "Integration — เชื่อมระบบภายนอก", size=32, color=WHITE, bold=True)
add_text(s, Inches(0.5), Inches(0.9), Inches(12), Inches(0.5),
         "Core DB รับข้อมูลจาก 3 แหล่ง ผ่าน Integration API", size=18, color=ACCENT_GOLD)

# Left: External sources
sources = [
    ("Space System", "ทีมอื่น", [
        "ข้อมูลโซน พื้นที่ บูท",
        "Floor Plan + Layout",
        "สถานะพื้นที่ว่าง/ไม่ว่าง",
        "ราคาตาม Tier",
    ], ACCENT_BLUE),
    ("IoT API", "ทีมอื่น", [
        "People Counting (sensors)",
        "Heatmap ความหนาแน่น",
        "Movement Flow / Dwell Time",
        "Zone-level analytics",
    ], ACCENT_ORANGE),
    ("SRT API", "การรถไฟ (ลูกค้า)", [
        "ตารางรถไฟ real-time",
        "ข้อมูลสถานี",
        "ระบบ Booking",
        "Migrate User จากแอปเดิม",
    ], ACCENT_GREEN),
]

for i, (name, team, items, color) in enumerate(sources):
    y = Inches(1.6 + i * 1.9)
    add_shape(s, Inches(0.5), y, Inches(5), Inches(1.7), border_color=color)
    add_text(s, Inches(0.7), y + Inches(0.05), Inches(2.2), Inches(0.4),
             f"{name} ({team})", size=14, color=color, bold=True)
    add_bullets(s, Inches(2.8), y + Inches(0.05), Inches(2.5), Inches(1.5),
                items, size=11, color=LIGHT_GRAY, spacing=4)

# Center: Arrow
add_text(s, Inches(5.5), Inches(3.5), Inches(1), Inches(0.5),
         "→", size=36, color=ACCENT_GOLD, bold=True, align=PP_ALIGN.CENTER)

# Center: Integration API + Core DB
add_shape(s, Inches(6.2), Inches(1.8), Inches(2.5), Inches(1.5),
          fill_color=RGBColor(0x1e, 0x1e, 0x30), border_color=ACCENT_GOLD, border_width=2)
add_text(s, Inches(6.2), Inches(1.9), Inches(2.5), Inches(0.5),
         "Integration API", size=15, color=ACCENT_GOLD, bold=True, align=PP_ALIGN.CENTER)
add_text(s, Inches(6.2), Inches(2.4), Inches(2.5), Inches(0.6),
         "แปลงข้อมูล\nให้เข้ากับระบบเรา", size=12, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

add_text(s, Inches(7.1), Inches(3.3), Inches(0.8), Inches(0.5),
         "↓", size=30, color=ACCENT_GOLD, bold=True, align=PP_ALIGN.CENTER)

add_shape(s, Inches(6.2), Inches(3.8), Inches(2.5), Inches(1.5),
          fill_color=RGBColor(0x2a, 0x25, 0x00), border_color=ACCENT_GOLD, border_width=3)
add_text(s, Inches(6.2), Inches(3.9), Inches(2.5), Inches(0.5),
         "Core Database", size=16, color=ACCENT_GOLD, bold=True, align=PP_ALIGN.CENTER)
add_text(s, Inches(6.2), Inches(4.4), Inches(2.5), Inches(0.6),
         "รวมข้อมูล\nSingle Source of Truth", size=12, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# Right: Arrow to outputs
add_text(s, Inches(8.7), Inches(3.5), Inches(1), Inches(0.5),
         "→", size=36, color=ACCENT_GREEN, bold=True, align=PP_ALIGN.CENTER)

# Right: NDF outputs
outputs = [
    ("Core Dashboard", "Report, Analytics", ACCENT_GOLD),
    ("ตู้ Kiosk", "แสดงผลข้อมูล", ACCENT_ORANGE),
    ("Mobile App", "ให้บริการ User", ACCENT_PURPLE),
    ("Admin Dashboard", "จัดการระบบ", ACCENT_BLUE),
]
for i, (name, desc, color) in enumerate(outputs):
    y = Inches(1.6 + i * 1.4)
    add_shape(s, Inches(9.5), y, Inches(3.5), Inches(1.1), border_color=color)
    add_text(s, Inches(9.7), y + Inches(0.05), Inches(3.1), Inches(0.4),
             name, size=14, color=color, bold=True)
    add_text(s, Inches(9.7), y + Inches(0.5), Inches(3.1), Inches(0.4),
             desc, size=12, color=LIGHT_GRAY)


# ============================================================
# SLIDE 17 — Roadmap
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_text(s, Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
         "Roadmap — แผนการพัฒนา", size=32, color=WHITE, bold=True)
add_text(s, Inches(0.5), Inches(0.9), Inches(12), Inches(0.5),
         "แบ่งเป็น 2 Phase — ส่งมอบทีละส่วน มั่นใจได้ว่าทำได้จริง", size=18, color=ACCENT_GOLD)

# Phase 1
add_shape(s, Inches(0.5), Inches(1.6), Inches(5.8), Inches(5.2),
          border_color=ACCENT_GOLD, border_width=3)
add_text(s, Inches(0.5), Inches(1.7), Inches(5.8), Inches(0.5),
         "Phase 1 — MVP (ทำแน่)", size=22, color=ACCENT_GOLD, bold=True, align=PP_ALIGN.CENTER)
add_bullets(s, Inches(0.8), Inches(2.4), Inches(5.2), Inches(4), [
    "Core Dashboard + Database",
    "Auth กลาง (Login ทุกระบบ)",
    "User & Role Management",
    "ตู้ Kiosk — AI, Route Finder, Train Schedule",
    "Mobile App — Login, Schedule, Booking, Navigation",
    "Admin Dashboard — Kiosk Management, Report",
    "Integration API — เชื่อม Space + IoT + SRT",
    "Report & Export (PDF/Excel)",
], size=15, color=LIGHT_GRAY, spacing=10)

# Phase 2
add_shape(s, Inches(6.8), Inches(1.6), Inches(5.8), Inches(5.2),
          border_color=DIM_GRAY, border_width=2)
add_text(s, Inches(6.8), Inches(1.7), Inches(5.8), Inches(0.5),
         "Phase 2 — Enhancement (ทำภายหลัง)", size=22, color=DIM_GRAY, bold=True, align=PP_ALIGN.CENTER)
add_bullets(s, Inches(7.1), Inches(2.4), Inches(5.2), Inches(4), [
    "Point / Reward System",
    "Campaign Management",
    "Advanced Analytics & AI Insights",
    "Advertisement Management",
    "Payment / Refund ในแอป",
    "Mini-ChatBot (Mobile)",
    "Revenue Forecast",
    "Multi-language expansion",
], size=15, color=DIM_GRAY, spacing=10)


# ============================================================
# SLIDE 18 — ปิด
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)

add_text(s, Inches(1), Inches(1.5), Inches(11), Inches(1),
         "NDF", size=60, color=ACCENT_GOLD, bold=True, align=PP_ALIGN.CENTER)
add_text(s, Inches(1), Inches(2.8), Inches(11), Inches(0.8),
         "Smart Tourism — Core System", size=32, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

add_shape(s, Inches(4.5), Inches(3.8), Inches(4.3), Pt(4), fill_color=ACCENT_GOLD)

add_text(s, Inches(1), Inches(4.2), Inches(11), Inches(0.5),
         "พร้อมส่งมอบ — ทุกฟังก์ชัน ทุก User Group", size=20, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# Summary boxes
summary = [
    ("4", "Frontend\nSystems", ACCENT_ORANGE),
    ("2", "Backend\nAPIs", ACCENT_BLUE),
    ("1", "Auth กลาง\nทุกระบบ", ACCENT_RED),
    ("1", "Core\nDatabase", ACCENT_GOLD),
]
for i, (num, label, color) in enumerate(summary):
    x = Inches(1.5 + i * 2.8)
    add_shape(s, x, Inches(5.0), Inches(2.2), Inches(1.5), border_color=color, border_width=2)
    add_text(s, x, Inches(5.1), Inches(2.2), Inches(0.7),
             num, size=36, color=color, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, x, Inches(5.7), Inches(2.2), Inches(0.7),
             label, size=13, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

add_text(s, Inches(1), Inches(6.8), Inches(11), Inches(0.5),
         "NDF Co., Ltd.  |  Smart Tourism Solution", size=14, color=DIM_GRAY, align=PP_ALIGN.CENTER)


# ============================================================
# SAVE
# ============================================================
prs.save(OUTPUT_PATH)
print(f"Saved: {OUTPUT_PATH}")
print(f"Slides: {len(prs.slides)}")
